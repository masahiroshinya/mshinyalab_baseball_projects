from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


OUT = Path(__file__).resolve().parents[1] / "draft_ai" / "20260708_AI作成_ゼミ発表.pptx"

BLUE = RGBColor(31, 78, 121)
DARK = RGBColor(35, 35, 35)
GRAY = RGBColor(105, 105, 105)
LIGHT_BLUE = RGBColor(220, 235, 247)


def set_run(run, text, size=22, bold=False, color=DARK):
    run.text = text
    run.font.name = "Yu Gothic"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_footer(slide, page):
    box = slide.shapes.add_textbox(Inches(0.55), Inches(7.05), Inches(12.2), Inches(0.25))
    p = box.text_frame.paragraphs[0]
    p.text = f"打者のマーカーレス3次元モーションキャプチャーシステム開発 | {page}"
    p.font.name = "Yu Gothic"
    p.font.size = Pt(9)
    p.font.color.rgb = GRAY
    p.alignment = PP_ALIGN.RIGHT


def add_title(slide, title, subtitle=None):
    box = slide.shapes.add_textbox(Inches(0.55), Inches(0.35), Inches(12.2), Inches(0.7))
    p = box.text_frame.paragraphs[0]
    set_run(p.add_run(), title, 28, True, BLUE)
    line = slide.shapes.add_shape(1, Inches(0.55), Inches(1.12), Inches(12.2), Inches(0.02))
    line.fill.solid()
    line.fill.fore_color.rgb = BLUE
    line.line.color.rgb = BLUE
    if subtitle:
        sub = slide.shapes.add_textbox(Inches(0.6), Inches(1.18), Inches(11.8), Inches(0.35))
        p = sub.text_frame.paragraphs[0]
        p.text = subtitle
        p.font.name = "Yu Gothic"
        p.font.size = Pt(13)
        p.font.color.rgb = GRAY


def add_bullets(slide, bullets, left=0.8, top=1.55, width=11.7, height=4.9, size=21):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    for i, text in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.level = 0
        p.font.name = "Yu Gothic"
        p.font.size = Pt(size)
        p.font.color.rgb = DARK
        p.space_after = Pt(9)


def add_table(slide, rows, left, top, width, height, font_size=13):
    table = slide.shapes.add_table(len(rows), len(rows[0]), Inches(left), Inches(top), Inches(width), Inches(height)).table
    for row_i, row in enumerate(rows):
        for col_i, value in enumerate(row):
            cell = table.cell(row_i, col_i)
            cell.text = value
            cell.margin_left = Inches(0.06)
            cell.margin_right = Inches(0.06)
            cell.fill.solid()
            cell.fill.fore_color.rgb = LIGHT_BLUE if row_i == 0 else RGBColor(255, 255, 255)
            for p in cell.text_frame.paragraphs:
                p.font.name = "Yu Gothic"
                p.font.size = Pt(font_size)
                p.font.bold = row_i == 0
                p.font.color.rgb = BLUE if row_i == 0 else DARK
    return table


def add_process(slide, steps):
    left = 0.75
    top = 1.55
    width = 11.85
    step_h = 0.63
    for i, step in enumerate(steps, start=1):
        y = top + (i - 1) * 0.82
        num = slide.shapes.add_shape(1, Inches(left), Inches(y), Inches(0.55), Inches(step_h))
        num.fill.solid()
        num.fill.fore_color.rgb = BLUE
        num.line.color.rgb = BLUE
        p = num.text_frame.paragraphs[0]
        p.text = str(i)
        p.font.name = "Yu Gothic"
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER
        box = slide.shapes.add_textbox(Inches(left + 0.72), Inches(y + 0.03), Inches(width - 0.72), Inches(step_h))
        p = box.text_frame.paragraphs[0]
        p.text = step
        p.font.name = "Yu Gothic"
        p.font.size = Pt(19)
        p.font.color.rgb = DARK


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    slide = prs.slides.add_slide(blank)
    box = slide.shapes.add_textbox(Inches(0.8), Inches(1.35), Inches(11.8), Inches(1.5))
    p = box.text_frame.paragraphs[0]
    set_run(p.add_run(), "打者のマーカーレス3次元\nモーションキャプチャーシステム開発", 34, True, BLUE)
    sub = slide.shapes.add_textbox(Inches(0.85), Inches(3.5), Inches(10.9), Inches(0.8))
    p = sub.text_frame.paragraphs[0]
    p.text = "カメラ映像のみから身体動作とバット軌道を3次元再構成する"
    p.font.name = "Yu Gothic"
    p.font.size = Pt(22)
    p.font.color.rgb = DARK
    meta = slide.shapes.add_textbox(Inches(0.9), Inches(5.85), Inches(8.0), Inches(0.5))
    p = meta.text_frame.paragraphs[0]
    p.text = "山藤 | 進矢研究室 | 2026-07-08"
    p.font.name = "Yu Gothic"
    p.font.size = Pt(16)
    p.font.color.rgb = GRAY
    add_footer(slide, 1)

    bullet_slides = [
        (
            "背景",
            [
                "野球の打撃動作は、骨盤・体幹・上肢・バットが短時間で連鎖的に動く高速運動である",
                "反射マーカー式モーションキャプチャーは高精度だが、高価で実験室環境に制限されやすい",
                "市販アクションカメラとAI姿勢推定を使えば、現場に近い環境で3D計測できる可能性がある",
            ],
        ),
        (
            "研究に対するモチベーション",
            [
                "高価な専用設備なしに、選手の自然なスイングを記録したい",
                "身体の動きだけでなく、バットの軌道も同時に扱いたい",
                "将来的には、バット速度、身体各部の角速度、運動連鎖のタイミングを定量化したい",
            ],
        ),
    ]
    for i, (title, bullets) in enumerate(bullet_slides, start=2):
        slide = prs.slides.add_slide(blank)
        add_title(slide, title)
        add_bullets(slide, bullets)
        add_footer(slide, i)

    slide = prs.slides.add_slide(blank)
    add_title(slide, "先行研究", "同じ方向に向いた研究と、本研究で引き継ぐ視点")
    add_table(
        slide,
        [
            ["方向性", "例", "本研究との関係"],
            ["打撃の3D解析", "Welch et al. など", "骨盤・肩・腕・バットの運動連鎖を定量化する視点"],
            ["条件別の打撃解析", "Fortenbaugh など", "球種・コースごとの動作差を分析する視点"],
            ["マーカーレス姿勢推定", "OpenPose, YOLO-pose など", "映像から身体キーポイントを取得する技術基盤"],
            ["簡易マーカーレス計測", "OpenCap など", "低コスト化の方向性。ただし高速スイングへの適用は課題"],
        ],
        0.55,
        1.45,
        12.2,
        4.65,
        12,
    )
    add_footer(slide, 4)

    slide = prs.slides.add_slide(blank)
    add_title(slide, "まだ解決されていないこと")
    add_bullets(
        slide,
        [
            "高速スイングでは、1フレームの同期ずれでも3D位置や速度推定に影響する",
            "屋外撮影では、カメラ設置、手ブレ補正、キャリブレーションの再現性が問題になる",
            "身体姿勢だけでなく、細長く高速に動くバットの追跡が難しい",
            "2D検出、同期、キャリブレーション、三角測量、可視化を一貫して扱うワークフローが必要",
        ],
        size=20,
    )
    add_footer(slide, 5)

    slide = prs.slides.add_slide(blank)
    add_title(slide, "どうやったら解決できるか")
    add_process(
        slide,
        [
            "DJI OSMO Action 3で 1080p / 240 fps の動画を撮影する",
            "手拍子音を用いた音声クロスコリレーションで動画を同期する",
            "ChArUcoボードでステレオキャリブレーションを行う",
            "YOLOv8-poseで身体17点、YOLOv8系検出器でバットを検出する",
            "三角測量で3D座標を算出し、スティックフィギュアとバット軌道を可視化する",
        ],
    )
    add_footer(slide, 6)

    slide = prs.slides.add_slide(blank)
    add_title(slide, "目的", "どの数値がどの程度になったら達成とみなすか")
    add_table(
        slide,
        [
            ["評価項目", "暫定的な到達目標"],
            ["動画同期精度", "1フレーム未満、240 fps換算で約4.17 ms未満"],
            ["ステレオキャリブレーション", "再投影誤差 RMS 2 px以下を目標"],
            ["身体キーポイント検出", "主要関節の有効検出率 90%以上を目標"],
            ["バット追跡", "スイング主要区間で軌道を連続的に描画できること"],
            ["3D再構成の実用精度", "既知寸法と比較し、数 cmオーダーの誤差に収めることを目標"],
        ],
        0.55,
        1.45,
        12.2,
        4.8,
        12,
    )
    add_footer(slide, 7)

    slide = prs.slides.add_slide(blank)
    add_title(slide, "評価方法")
    add_bullets(
        slide,
        [
            "同期: 手拍子音のピーク位置と相互相関のずれを確認する",
            "キャリブレーション: 再投影誤差 RMS と、複数回撮影でのばらつきを確認する",
            "3D再構成: 肩幅、身長、バット長など既知寸法と比較する",
            "バット軌道: 手首位置との整合性、欠損フレーム数、軌道の連続性を確認する",
            "実用性: 撮影から可視化までの手順が再現可能かを確認する",
        ],
        size=19,
    )
    add_footer(slide, 8)

    slide = prs.slides.add_slide(blank)
    add_title(slide, "今後の進め方")
    add_bullets(
        slide,
        [
            "実撮影データで、同期・検出・3D再構成を一貫して検証する",
            "数値目標を実測値で更新する",
            "結果画像、3D可視化、バット軌道の図を発表資料に追加する",
            "最終版は 人間作製/final に保存する",
        ],
    )
    add_footer(slide, 9)

    prs.save(OUT)
    print(OUT)


if __name__ == "__main__":
    main()
